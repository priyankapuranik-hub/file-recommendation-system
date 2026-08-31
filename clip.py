import os
import hashlib
import numpy as np
from transformers import CLIPModel, CLIPProcessor
from typing import List, Dict, Tuple, Union
import faiss
import mimetypes
import re
from PIL import Image
import torch

class FileRecommendationSystem:
    def __init__(self, base_path: str, model_name: str = 'openai/clip-vit-large-patch14'):
        """
        Initialize the File Recommendation System with CLIP model for both text and image understanding
        
        Args:
            base_path (str): Root directory to scan for files
            model_name (str): CLIP model name (e.g., 'openai/clip-vit-large-patch14')
        """
        self.base_path = base_path
        self.model_name = model_name
        self.model = CLIPModel.from_pretrained(model_name)
        self.processor = CLIPProcessor.from_pretrained(model_name)
        
        # System and temporary file extensions to exclude
        self.excluded_extensions = {
            '.tmp', '.log', '.swp', '.lock', 
            '.sys', '.dat', '.ini', '.config',
            '.exe', '.dll', '.bin', 
            '.DS_Store', 'desktop.ini', 'thumbs.db'
        }
        
        # Supported file extensions
        self.supported_extensions = {'.txt', '.pdf', '.docx', '.ppt', '.pptx', '.jpg', '.jpeg', '.xlsx', '.xls', '.png', '.cpp', '.c', '.py', '.java', '.js', '.html', '.css'}
        
        # Vector database using FAISS
        self.dimension = 768  # CLIP-vit-large-patch14 embedding size
        self.index = faiss.IndexFlatL2(self.dimension)
        
        # Metadata storage
        self.file_metadata = {}
        
        # Check if model is CLIP-based
        self.is_clip_model = 'clip' in model_name.lower()
        print(f"Using {'CLIP-based' if self.is_clip_model else 'unknown'} model: {model_name}")

    def is_valid_file(self, file_path: str) -> bool:
        """
        Check if file should be included in recommendation system
        
        Args:
            file_path (str): Path to the file
        
        Returns:
            bool: Whether file is valid for processing
        """
        filename = os.path.basename(file_path)
        file_ext = os.path.splitext(filename)[1].lower()
        
        conditions = [
            not filename.startswith('.'),  # Exclude hidden files
            file_ext not in self.excluded_extensions,
            file_ext in self.supported_extensions,  # Only allow supported files
            os.path.isfile(file_path),  # Ensure it's a file
            os.path.getsize(file_path) > 0  # Exclude empty files
        ]
        
        return all(conditions)
    
    def clean_text(self, text: str) -> str:
        """
        Clean extracted text by removing excessive whitespace and special characters
        
        Args:
            text (str): Raw extracted text
        
        Returns:
            str: Cleaned text
        """
        text = re.sub(r'\s+', ' ', text.strip())
        return text if text else "No text extracted"
    
    def is_image_file(self, file_path: str) -> bool:
        """
        Check if file is an image
        
        Args:
            file_path (str): Path to the file
        
        Returns:
            bool: Whether file is an image
        """
        mime_type, _ = mimetypes.guess_type(file_path)
        file_ext = os.path.splitext(file_path)[1].lower()
        return (mime_type in ['image/jpeg', 'image/png'] or 
                file_ext in ['.jpg', '.jpeg', '.png'])
    
    def process_image(self, file_path: str) -> Tuple[str, Union[np.ndarray, None]]:
        """
        Process image file using both OCR and CLIP visual understanding
        
        Args:
            file_path (str): Path to the image file
            
        Returns:
            Tuple[str, Union[np.ndarray, None]]: OCR text and image embedding
        """
        try:
            import pytesseract
            # Load image
            image = Image.open(file_path).convert('RGB')
            
            # Extract text via OCR
            ocr_text = pytesseract.image_to_string(image)
            ocr_text = self.clean_text(ocr_text)
            
            # Get visual embedding if using a CLIP model
            if self.is_clip_model:
                # Preprocess image for CLIP
                inputs = self.processor(images=image, return_tensors="pt")
                with torch.no_grad():
                    image_features = self.model.get_image_features(**inputs)
                    image_embedding = image_features.cpu().numpy()[0]
                return ocr_text, image_embedding
            else:
                # For non-CLIP models, just return OCR text
                return ocr_text, None
                
        except ImportError:
            return f"Error processing image {os.path.basename(file_path)}: pytesseract module not installed", None
        except Exception as e:
            return f"Error processing image {os.path.basename(file_path)}: {str(e)}", None
    
    def extract_file_text(self, file_path: str) -> Union[str, Tuple[str, Union[np.ndarray, None]]]:
        """
        Extract text content from different file types
        
        Args:
            file_path (str): Path to the file
        
        Returns:
            Union[str, Tuple[str, Union[np.ndarray, None]]]: 
                For text files: extracted text content
                For images: (text content, image embedding)
        """
        try:
            mime_type, _ = mimetypes.guess_type(file_path)
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # Handle images using CLIP model for visual understanding
            if self.is_image_file(file_path):
                return self.process_image(file_path)
            
            # Text files
            if mime_type and mime_type.startswith('text/'):
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return self.clean_text(f.read())
            
            # PDF files
            elif mime_type == 'application/pdf' or file_ext == '.pdf':
                try:
                    import PyPDF2
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: PyPDF2 module not installed"
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ' '.join([page.extract_text() or "" for page in reader.pages])
                    return self.clean_text(text)
            
            # DOCX files
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or file_ext == '.docx':
                try:
                    import docx
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: python-docx module not installed"
                doc = docx.Document(file_path)
                text = ' '.join([para.text for para in doc.paragraphs if para.text])
                return self.clean_text(text)
            
            # PPT/PPTX files
            elif mime_type in ['application/vnd.ms-powerpoint', 
                             'application/vnd.openxmlformats-officedocument.presentationml.presentation'] or file_ext in ['.ppt', '.pptx']:
                try:
                    from pptx import Presentation
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: python-pptx module not installed"
                prs = Presentation(file_path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            text_parts.append(shape.text)
                text = ' '.join(text_parts)
                return self.clean_text(text)
            
            # Excel files (XLSX/XLS)
            elif mime_type in ['application/vnd.ms-excel', 
                             'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'] or file_ext in ['.xlsx', '.xls']:
                try:
                    import pandas as pd
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: pandas module not installed"
                try:
                    xls = pd.ExcelFile(file_path)
                    text_parts = []
                    for sheet_name in xls.sheet_names:
                        df = pd.read_excel(file_path, sheet_name=sheet_name)
                        text_parts.append(f"Sheet: {sheet_name}")
                        text_parts.append(str(df.head(100)))  # Limit to first 100 rows
                    text = "\n".join(text_parts)
                    return self.clean_text(text)
                except Exception as e:
                    return f"Error extracting text from {os.path.basename(file_path)}: {str(e)}"
            
            return f"File type {mime_type or file_ext} not supported for text extraction"
        
        except Exception as e:
            return f"Error extracting text from {os.path.basename(file_path)}: {str(e)}"
    
    def generate_file_embeddings(self):
        """
        Scan directory, generate embeddings for valid files and store in vector database
        """
        if not os.path.exists(self.base_path):
            print(f"Directory {self.base_path} does not exist. No files processed.")
            return
        
        # Reset FAISS index and metadata to ensure fresh processing
        self.index = faiss.IndexFlatL2(self.dimension)
        self.file_metadata = {}
        
        print(f"Using {self.model_name} with dimension {self.dimension}")
        print(f"CLIP-based model: {self.is_clip_model}")
        
        for root, _, files in os.walk(self.base_path):
            for filename in files:
                file_path = os.path.join(root, filename)
                
                if self.is_valid_file(file_path):
                    try:
                        # Generate a unique hash for the file
                        file_hash = hashlib.md5(file_path.encode()).hexdigest()
                        mod_time = os.path.getmtime(file_path)
                        
                        # Extract text content and/or image embedding
                        result = self.extract_file_text(file_path)
                        
                        # Handle different return types based on file type
                        if isinstance(result, tuple) and len(result) == 2:
                            # Image file with text and embedding
                            file_text, image_embedding = result
                            
                            if file_text.startswith("Error"):
                                print(file_text)
                                continue
                                
                            # Use image embedding if available, otherwise encode OCR text
                            if image_embedding is not None:
                                embedding = image_embedding
                            else:
                                # Encode OCR text with CLIP text encoder
                                inputs = self.processor(text=file_text, return_tensors="pt", padding=True, truncation=True)
                                with torch.no_grad():
                                    embedding = self.model.get_text_features(**inputs).cpu().numpy()[0]
                        else:
                            # Text-only file
                            file_text = result
                            
                            # Skip files with extraction errors
                            if file_text.startswith("Error"):
                                print(file_text)
                                continue
                            
                            # Encode text with CLIP text encoder
                            inputs = self.processor(text=file_text, return_tensors="pt", padding=True, truncation=True)
                            with torch.no_grad():
                                embedding = self.model.get_text_features(**inputs).cpu().numpy()[0]
                        
                        # Add to FAISS index
                        self.index.add(np.array([embedding]))
                        
                        # Store metadata
                        self.file_metadata[file_hash] = {
                            'path': file_path,
                            'filename': filename,
                            'size': os.path.getsize(file_path),
                            'mod_time': mod_time,
                            'is_image': self.is_image_file(file_path)
                        }
                        print(f"Processed {filename}")
                    
                    except Exception as e:
                        print(f"Error processing {file_path}: {e}")
    
    def recommend_files(self, query: str, top_k: int = 5, image_path: str = None) -> List[Dict]:
        """
        Recommend files based on text query or image query
        
        Args:
            query (str): User's text query
            top_k (int): Number of top recommendations to return
            image_path (str, optional): Path to query image
            
        Returns:
            List of recommended files with metadata
        """
        if not self.file_metadata:
            print("No files indexed. Run generate_file_embeddings first.")
            return []
        
        # Generate embedding for query
        if image_path and self.is_clip_model and os.path.exists(image_path):
            # Query by image
            try:
                image = Image.open(image_path).convert('RGB')
                inputs = self.processor(images=image, return_tensors="pt")
                with torch.no_grad():
                    query_embedding = self.model.get_image_features(**inputs).cpu().numpy()[0]
            except Exception as e:
                print(f"Error processing query image: {e}")
                # Fall back to text query if image processing fails
                inputs = self.processor(text=query, return_tensors="pt", padding=True, truncation=True)
                with torch.no_grad():
                    query_embedding = self.model.get_text_features(**inputs).cpu().numpy()[0]
        else:
            # Query by text
            inputs = self.processor(text=query, return_tensors="pt", padding=True, truncation=True)
            with torch.no_grad():
                query_embedding = self.model.get_text_features(**inputs).cpu().numpy()[0]
        
        # Search in vector database
        distances, indices = self.index.search(np.array([query_embedding]), min(top_k, len(self.file_metadata)))
        
        # Retrieve and return recommended files
        recommendations = []
        for dist, idx in zip(distances[0], indices[0]):
            if idx < len(self.file_metadata):
                file_hash = list(self.file_metadata.keys())[idx]
                file_info = self.file_metadata[file_hash].copy()
                file_info['similarity_score'] = 1 / (1 + dist)  # Convert distance to similarity
                recommendations.append(file_info)
        
        return recommendations

# Example usage
if __name__ == "__main__":
    # Initialize the system with base path
    recommender = FileRecommendationSystem("E:\\c++ programs", model_name="E:/File_recommend_sys/clip-vit-large-patch14")
    
    # Print embedding model information
    print(f"Using embedding model: {recommender.model_name}")
    print(f"Embedding dimension: {recommender.dimension}")
    
    # Generate embeddings for all files
    recommender.generate_file_embeddings()
    
    # Example text recommendation
    results = recommender.recommend_files("area of circle")
    if results:
        print("\nRecommended files for text query:")
        for result in results:
            print(f"Recommended File: {result['filename']}")
            print(f"Path: {result['path']}")
            print(f"Similarity Score: {result['similarity_score']:.4f}")
            print(f"Is Image: {result['is_image']}\n")
    else:
        print("No recommendations found.")
    
    # Example image recommendation
    sample_image = "E:\\Documents\\sample.jpg"  # Replace with an actual image path
    if os.path.exists(sample_image):
        print("\nTesting image-based query:")
        image_results = recommender.recommend_files("pheromones", image_path=sample_image)
        if image_results:
            print("\nRecommended files for image query:")
            for result in image_results:
                print(f"Recommended File: {result['filename']}")
                print(f"Path: {result['path']}")
                print(f"Similarity Score: {result['similarity_score']:.4f}")
                print(f"Is Image: {result['is_image']}\n")
        else:
            print("No image-based recommendations found.")