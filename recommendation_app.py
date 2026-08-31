import os
import hashlib
import numpy as np
from sentence_transformers import SentenceTransformer
from typing import List, Dict
import faiss
import mimetypes
import re
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import ttkbootstrap as ttkb
from ttkbootstrap.constants import *
from tkinterdnd2 import *
import shutil
import threading
from pathlib import Path

# FileRecommendationSystem class (unchanged from previous version without cache)
class FileRecommendationSystem:
    def __init__(self, base_path: str, embedding_model: str = 'all-MiniLM-L6-v2'):
        self.base_path = base_path
        self.embedding_model = embedding_model
        self.model = SentenceTransformer(embedding_model)
        
        self.excluded_extensions = {
            '.tmp', '.log', '.swp', '.lock', 
            '.sys', '.dat', '.ini', '.config',
            '.exe', '.dll', '.bin', 
            '.DS_Store', 'desktop.ini', 'thumbs.db'
        }
        
        self.supported_extensions = {'.txt', '.pdf', '.docx', '.ppt', '.pptx', '.jpg', '.jpeg', '.xlsx', '.xls', '.png'}
        
        self.dimension = self.model.get_sentence_embedding_dimension()
        self.index = faiss.IndexFlatL2(self.dimension)
        self.file_metadata = {}

    def is_valid_file(self, file_path: str) -> bool:
        filename = os.path.basename(file_path)
        file_ext = os.path.splitext(filename)[1].lower()
        
        conditions = [
            not filename.startswith('.'),
            file_ext not in self.excluded_extensions,
            file_ext in self.supported_extensions,
            os.path.isfile(file_path),
            os.path.getsize(file_path) > 0
        ]
        
        return all(conditions)
    
    def clean_text(self, text: str) -> str:
        text = re.sub(r'\s+', ' ', text.strip())
        return text if text else "No text extracted"
    
    def extract_file_text(self, file_path: str) -> str:
        try:
            mime_type, _ = mimetypes.guess_type(file_path)
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if mime_type and mime_type.startswith('text/'):
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return self.clean_text(f.read())
            
            elif mime_type == 'application/pdf' or file_ext == '.pdf':
                try:
                    import PyPDF2
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: PyPDF2 module not installed"
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ' '.join([page.extract_text() or "" for page in reader.pages])
                    return self.clean_text(text)
            
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or file_ext == '.docx':
                try:
                    import docx
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: python-docx module not installed"
                doc = docx.Document(file_path)
                text = ' '.join([para.text for para in doc.paragraphs if para.text])
                return self.clean_text(text)
            
            elif mime_type in ['application/vnd.ms-powerpoint', 
                             'application/vnd.openxmlformats-officedocument.presentationml.presentation'] or file_ext in ['.ppt', '.pptx']:
                try:
                    from pptx import Presentation
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: python-pptx module not installed"
                prs = Presentation(file_path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            text_parts.append(shape.text)
                text = ' '.join(text_parts)
                return self.clean_text(text)
            
            elif mime_type in ['application/vnd.ms-excel', 
                             'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'] or file_ext in ['.xlsx', '.xls']:
                try:
                    import pandas as pd
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: pandas module not installed"
                try:
                    xls = pd.ExcelFile(file_path)
                    text_parts = []
                    for sheet_name in xls.sheet_names:
                        df = pd.read_excel(file_path, sheet_name=sheet_name)
                        text_parts.append(f"Sheet: {sheet_name}")
                        text_parts.append(str(df.head(100)))
                    text = "\n".join(text_parts)
                    return self.clean_text(text)
                except Exception as e:
                    return f"Error extracting text from {os.path.basename(file_path)}: {str(e)}"
            
            elif mime_type in ['image/jpeg', 'image/png'] or file_ext in ['.jpg', '.jpeg', '.png']:
                try:
                    from PIL import Image
                    import pytesseract
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: PIL or pytesseract module not installed"
                image = Image.open(file_path)
                text = pytesseract.image_to_string(image)
                return self.clean_text(text)
            
            return f"File type {mime_type or file_ext} not supported for text extraction"
        
        except Exception as e:
            return f"Error extracting text from {os.path.basename(file_path)}: {str(e)}"
    
    def generate_file_embeddings(self):
        if not os.path.exists(self.base_path):
            print(f"Directory {self.base_path} does not exist. No files processed.")
            return
        
        self.index = faiss.IndexFlatL2(self.dimension)
        self.file_metadata = {}
        
        for root, _, files in os.walk(self.base_path):
            for filename in files:
                file_path = os.path.join(root, filename)
                
                if self.is_valid_file(file_path):
                    try:
                        file_hash = hashlib.md5(file_path.encode()).hexdigest()
                        mod_time = os.path.getmtime(file_path)
                        
                        file_text = self.extract_file_text(file_path)
                        
                        if file_text.startswith("Error extracting"):
                            print(file_text)
                            continue
                        
                        embedding = self.model.encode(file_text)
                        self.index.add(np.array([embedding]))
                        
                        self.file_metadata[file_hash] = {
                            'path': file_path,
                            'filename': filename,
                            'size': os.path.getsize(file_path),
                            'mod_time': mod_time
                        }
                        print(f"Processed {filename}")
                    
                    except Exception as e:
                        print(f"Error processing {file_path}: {e}")
    
    def recommend_files(self, query: str, top_k: int = 5) -> List[Dict]:
        if not self.file_metadata:
            print("No files indexed. Run generate_file_embeddings first.")
            return []
        
        query_embedding = self.model.encode(query)
        distances, indices = self.index.search(np.array([query_embedding]), min(top_k, len(self.file_metadata)))
        
        recommendations = []
        for dist, idx in zip(distances[0], indices[0]):
            file_hash = list(self.file_metadata.keys())[idx]
            file_info = self.file_metadata[file_hash].copy()
            file_info['similarity_score'] = 1 / (1 + dist)
            recommendations.append(file_info)
        
        return recommendations

# Tkinter-based UI Application
class FileRecommendationApp:
    def __init__(self, root):
        self.root = root
        self.root.title("File Recommendation System")
        self.root.geometry("800x600")
        self.recommender = FileRecommendationSystem("E:\\Documents")
        
        # Configure style
        self.style = ttkb.Style(theme='minty')
        
        # Main frame
        self.main_frame = ttk.Frame(self.root, padding=10)
        self.main_frame.pack(fill=BOTH, expand=True)
        
        # Drag-and-drop area
        self.drop_frame = ttk.LabelFrame(self.main_frame, text="Drag & Drop Files Here", padding=10)
        self.drop_frame.pack(fill=X, pady=5)
        self.drop_label = ttk.Label(self.drop_frame, text="Drop files or click to browse", anchor="center", font=("Helvetica", 12))
        self.drop_label.pack(fill=BOTH, expand=True)
        
        # Bind drag-and-drop and click events
        try:
            self.drop_label.drop_target_register(DND_FILES)
            self.drop_label.dnd_bind('<<Drop>>', self.handle_drop)
        except Exception as e:
            self.drop_label.configure(text="Drop files or click to browse (Drag-and-drop not available)")
            print(f"Drag-and-drop setup failed: {e}")
        self.drop_label.bind("<Button-1>", self.browse_files)
        
        # Process button and progress bar
        self.process_frame = ttk.Frame(self.main_frame)
        self.process_frame.pack(fill=X, pady=5)
        self.process_button = ttk.Button(self.process_frame, text="Process Files", command=self.start_processing, style='primary.TButton')
        self.process_button.pack(side=LEFT)
        self.progress = ttk.Progressbar(self.process_frame, mode='determinate', length=200)
        self.progress.pack(side=LEFT, padx=10)
        
        # Search bar
        self.search_frame = ttk.Frame(self.main_frame)
        self.search_frame.pack(fill=X, pady=5)
        ttk.Label(self.search_frame, text="Search Query:").pack(side=LEFT)
        self.search_entry = ttk.Entry(self.search_frame)
        self.search_entry.pack(side=LEFT, fill=X, expand=True, padx=5)
        self.search_button = ttk.Button(self.search_frame, text="Search", command=self.search_files, style='success.TButton')
        self.search_button.pack(side=LEFT)
        
        # Results table
        self.results_frame = ttk.LabelFrame(self.main_frame, text="Recommendations", padding=10)
        self.results_frame.pack(fill=BOTH, expand=True, pady=5)
        self.tree = ttk.Treeview(self.results_frame, columns=("Filename", "Path", "Score"), show='headings')
        self.tree.heading("Filename", text="Filename")
        self.tree.heading("Path", text="Path")
        self.tree.heading("Score", text="Similarity Score")
        self.tree.column("Filename", width=200)
        self.tree.column("Path", width=350)
        self.tree.column("Score", width=100)
        self.tree.pack(fill=BOTH, expand=True)
        self.tree.bind("<Double-1>", self.open_file)
        
        # Status bar
        self.status_var = tk.StringVar(value="Ready")
        self.status_bar = ttk.Label(self.main_frame, textvariable=self.status_var, relief=SUNKEN, anchor="w")
        self.status_bar.pack(fill=X, pady=5)
        
        # Initialize processing state
        self.processing = False
    
    def handle_drop(self, event):
        if self.processing:
            self.status_var.set("Cannot upload while processing")
            return
        files = self.root.splitlist(event.data)
        self.upload_files(files)
    
    def browse_files(self, event=None):
        if self.processing:
            self.status_var.set("Cannot upload while processing")
            return
        files = filedialog.askopenfilenames(
            filetypes=[("Supported Files", "*.txt *.pdf *.docx *.ppt *.pptx *.jpg *.jpeg *.xlsx *.xls *.png")]
        )
        if files:
            self.upload_files(files)
    
    def upload_files(self, files):
        supported_extensions = self.recommender.supported_extensions
        uploaded = 0
        for file in files:
            file_ext = os.path.splitext(file)[1].lower()
            if file_ext in supported_extensions:
                dest_path = os.path.join("E:\\Documents", os.path.basename(file))
                try:
                    shutil.copy(file, dest_path)
                    uploaded += 1
                except Exception as e:
                    self.status_var.set(f"Error uploading {os.path.basename(file)}: {str(e)}")
        self.status_var.set(f"Uploaded {uploaded} file(s) to E:\\Documents")
    
    def start_processing(self):
        if self.processing:
            self.status_var.set("Already processing")
            return
        self.processing = True
        self.process_button.configure(state='disabled')
        self.status_var.set("Processing files...")
        self.progress['value'] = 0
        
        # Run processing in a separate thread to avoid freezing UI
        threading.Thread(target=self.process_files, daemon=True).start()
    
    def process_files(self):
        try:
            files = []
            for root, _, filenames in os.walk("E:\\Documents"):
                for filename in filenames:
                    file_path = os.path.join(root, filename)
                    if self.recommender.is_valid_file(file_path):
                        files.append(file_path)
            
            total_files = len(files)
            if total_files == 0:
                self.root.after(0, self.processing_complete, "No valid files found")
                return
            
            self.recommender.generate_file_embeddings()
            
            # Simulate progress (since generate_file_embeddings logs progress)
            for i in range(total_files):
                self.progress['value'] = (i + 1) / total_files * 100
                self.root.update()
            
            self.root.after(0, self.processing_complete, f"Processed {total_files} file(s)")
        
        except Exception as e:
            self.root.after(0, self.processing_complete, f"Error processing files: {str(e)}")
    
    def processing_complete(self, message):
        self.status_var.set(message)
        self.process_button.configure(state='normal')
        self.processing = False
        self.progress['value'] = 100
    
    def search_files(self):
        query = self.search_entry.get().strip()
        if not query:
            self.status_var.set("Please enter a search query")
            return
        
        # Clear previous results
        for item in self.tree.get_children():
            self.tree.delete(item)
        
        try:
            results = self.recommender.recommend_files(query, top_k=5)
            if not results:
                self.status_var.set("No recommendations found")
                return
            
            for result in results:
                self.tree.insert("", END, values=(
                    result['filename'],
                    result['path'],
                    f"{result['similarity_score']:.4f}"
                ))
            self.status_var.set(f"Found {len(results)} recommendation(s)")
        
        except Exception as e:
            self.status_var.set(f"Error searching: {str(e)}")
    
    def open_file(self, event):
        item = self.tree.selection()
        if item:
            file_path = self.tree.item(item, "values")[1]
            try:
                os.startfile(file_path)
            except Exception as e:
                self.status_var.set(f"Error opening file: {str(e)}")

# Run the application
if __name__ == "__main__":
    root = TkinterDnD.Tk()
    app = FileRecommendationApp(root)
    root.mainloop()