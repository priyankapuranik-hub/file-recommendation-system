import os
import hashlib
import numpy as np
from sentence_transformers import SentenceTransformer
from typing import List, Dict
import faiss
import mimetypes
import re

class FileRecommendationSystem:
    def __init__(self, base_path: str, embedding_model: str = 'e5-large-v2'):
        self.base_path = base_path
        self.embedding_model = embedding_model
        self.model = SentenceTransformer(embedding_model)
        
        self.excluded_extensions = {
            '.tmp', '.log', '.swp', '.lock', 
            '.sys', '.dat', '.ini', '.config',
            '.exe', '.dll', '.bin', 
            '.DS_Store', 'desktop.ini', 'thumbs.db'
        }
        
        self.supported_extensions = {'.txt', '.pdf', '.docx', '.ppt', '.pptx', '.jpg', '.jpeg', '.xlsx', '.xls', '.png', '.cpp'}
        
        self.dimension = self.model.get_sentence_embedding_dimension()
        self.index = faiss.IndexFlatL2(self.dimension)
        self.file_metadata = {}

    def is_valid_file(self, file_path: str) -> bool:
        filename = os.path.basename(file_path)
        file_ext = os.path.splitext(filename)[1].lower()
        conditions = [
            not filename.startswith('.'),
            file_ext not in self.excluded_extensions,
            file_ext in self.supported_extensions,
            os.path.isfile(file_path),
            os.path.getsize(file_path) > 0
        ]
        return all(conditions)
    
    def clean_text(self, text: str) -> str:
        text = re.sub(r'\s+', ' ', text.strip())
        return text if text else "No text extracted"
    
    def extract_file_text(self, file_path: str) -> str:
        try:
            mime_type, _ = mimetypes.guess_type(file_path)
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if mime_type and mime_type.startswith('text/'):
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return self.clean_text(f.read())
            
            elif mime_type == 'application/pdf' or file_ext == '.pdf':
                import PyPDF2
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ' '.join([page.extract_text() or "" for page in reader.pages])
                    return self.clean_text(text)
            
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or file_ext == '.docx':
                import docx
                doc = docx.Document(file_path)
                text = ' '.join([para.text for para in doc.paragraphs if para.text])
                return self.clean_text(text)
            
            elif mime_type in ['application/vnd.ms-powerpoint', 
                               'application/vnd.openxmlformats-officedocument.presentationml.presentation'] or file_ext in ['.ppt', '.pptx']:
                from pptx import Presentation
                prs = Presentation(file_path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            text_parts.append(shape.text)
                return self.clean_text(' '.join(text_parts))
            
            elif mime_type in ['application/vnd.ms-excel', 
                               'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'] or file_ext in ['.xlsx', '.xls']:
                import pandas as pd
                xls = pd.ExcelFile(file_path)
                text_parts = []
                for sheet_name in xls.sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    text_parts.append(f"Sheet: {sheet_name}")
                    text_parts.append(str(df.head(100)))
                return self.clean_text("\n".join(text_parts))
            
            elif mime_type in ['image/jpeg', 'image/png'] or file_ext in ['.jpg', '.jpeg', '.png']:
                from PIL import Image
                import pytesseract
                image = Image.open(file_path)
                text = pytesseract.image_to_string(image)
                return self.clean_text(text)
            
            return f"File type {mime_type or file_ext} not supported for text extraction"
        
        except Exception as e:
            return f"Error extracting text from {os.path.basename(file_path)}: {str(e)}"
    
    def generate_file_embeddings(self):
        if not os.path.exists(self.base_path):
            print(f"Directory {self.base_path} does not exist.")
            return
        
        self.index = faiss.IndexFlatL2(self.dimension)
        self.file_metadata = {}
        
        for root, _, files in os.walk(self.base_path):
            for filename in files:
                file_path = os.path.join(root, filename)
                
                if self.is_valid_file(file_path):
                    try:
                        file_hash = hashlib.md5(file_path.encode()).hexdigest()
                        mod_time = os.path.getmtime(file_path)
                        file_text = self.extract_file_text(file_path)
                        
                        if file_text.startswith("Error extracting"):
                            print(file_text)
                            continue
                        
                        # 🔁 Add prefix for e5-style embedding
                        embedding = self.model.encode(f"passage: {file_text}")
                        
                        self.index.add(np.array([embedding]))
                        self.file_metadata[file_hash] = {
                            'path': file_path,
                            'filename': filename,
                            'size': os.path.getsize(file_path),
                            'mod_time': mod_time
                        }
                        print(f"Processed {filename}")
                    
                    except Exception as e:
                        print(f"Error processing {file_path}: {e}")
    
    def recommend_files(self, query: str, top_k: int = 5) -> List[Dict]:
        if not self.file_metadata:
            print("No files indexed. Run generate_file_embeddings first.")
            return []
        
        # 🔁 Add prefix for e5-style query
        query_embedding = self.model.encode(f"query: {query}")
        
        distances, indices = self.index.search(np.array([query_embedding]), min(top_k, len(self.file_metadata)))
        
        recommendations = []
        for dist, idx in zip(distances[0], indices[0]):
            file_hash = list(self.file_metadata.keys())[idx]
            file_info = self.file_metadata[file_hash].copy()
            file_info['similarity_score'] = 1 / (1 + dist)
            recommendations.append(file_info)
        
        return recommendations

# Example usage
if __name__ == "__main__":
    recommender = FileRecommendationSystem("E:\\Documents", embedding_model="e5-large-v2")
    recommender.generate_file_embeddings()
    
    results = recommender.recommend_files("cloud computing")
    if results:
        print("\nRecommended files:")
        for result in results:
            print(f"File: {result['filename']}")
            print(f"Path: {result['path']}")
            print(f"Similarity Score: {result['similarity_score']:.4f}\n")
    else:
        print("No recommendations found.")
