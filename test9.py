import os
import hashlib
import numpy as np
from sentence_transformers import SentenceTransformer
from typing import List, Dict
import faiss
import mimetypes
import re
from sklearn.preprocessing import normalize

class FileRecommendationSystem:
    def __init__(self, base_path: str, embedding_model: str = 'intfloat/e5-large-v2'):
        """
        Initialize the File Recommendation System
        
        Args:
            base_path (str): Root directory to scan for files
            embedding_model (str): Sentence Transformer model for generating embeddings
        """
        self.base_path = base_path
        self.embedding_model = embedding_model
        self.model = SentenceTransformer(embedding_model)
        
        # System and temporary file extensions to exclude
        self.excluded_extensions = {
            '.tmp', '.log', '.swp', '.lock', 
            '.sys', '.dat', '.ini', '.config',
            '.exe', '.dll', '.bin', 
            '.DS_Store', 'desktop.ini', 'thumbs.db'
        }
        
        # Supported file extensions
        self.supported_extensions = {'.txt', '.pdf', '.docx', '.ppt', '.pptx', '.jpg', '.jpeg', '.xlsx', '.xls', '.png', '.cpp'}
        
        # Vector database using FAISS
        self.dimension = self.model.get_sentence_embedding_dimension()
        self.index = faiss.IndexFlatL2(self.dimension)
        
        # Metadata storage
        self.file_metadata = {}

    def is_valid_file(self, file_path: str) -> bool:
        """
        Check if file should be included in recommendation system
        
        Args:
            file_path (str): Path to the file
        
        Returns:
            bool: Whether file is valid for processing
        """
        filename = os.path.basename(file_path)
        file_ext = os.path.splitext(filename)[1].lower()
        
        conditions = [
            not filename.startswith('.'),  # Exclude hidden files
            file_ext not in self.excluded_extensions,
            file_ext in self.supported_extensions,  # Only allow supported files
            os.path.isfile(file_path),  # Ensure it's a file
            os.path.getsize(file_path) > 0  # Exclude empty files
        ]
        
        return all(conditions)
    
    def clean_text(self, text: str) -> str:
        """
        Clean extracted text by removing excessive whitespace, special characters, and normalizing
        
        Args:
            text (str): Raw extracted text
        
        Returns:
            str: Cleaned text
        """
        text = re.sub(r'[^\w\s]', ' ', text.lower())  # Remove special chars, lowercase
        text = re.sub(r'\s+', ' ', text.strip())
        return text if text else "No text extracted"
    
    def chunk_text(self, text: str, max_length: int = 500) -> List[str]:
        """
        Split text into chunks to avoid token truncation
        
        Args:
            text (str): Input text
            max_length (int): Maximum length per chunk
        
        Returns:
            List[str]: List of text chunks
        """
        words = text.split()
        chunks = []
        current_chunk = []
        current_length = 0
        
        for word in words:
            current_length += len(word) + 1
            if current_length > max_length:
                chunks.append(' '.join(current_chunk))
                current_chunk = [word]
                current_length = len(word) + 1
            else:
                current_chunk.append(word)
        
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks if chunks else [text]
    
    def extract_file_text(self, file_path: str) -> str:
        """
        Extract text content from different file types
        
        Args:
            file_path (str): Path to the file
        
        Returns:
            str: Extracted text content or error message
        """
        try:
            mime_type, _ = mimetypes.guess_type(file_path)
            file_ext = os.path.splitext(file_path)[1].lower()
            filename = os.path.basename(file_path)
            
            # Text files (including .cpp)
            if mime_type and mime_type.startswith('text/') or file_ext == '.cpp':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
                    print(f"Extracted text from {filename}: {text[:100]}...")  # Debug
                    return self.clean_text(text)
            
            # PDF files
            elif mime_type == 'application/pdf' or file_ext == '.pdf':
                try:
                    import pdfplumber
                except ImportError:
                    return f"Error extracting text from {filename}: pdfplumber module not installed"
                with pdfplumber.open(file_path) as pdf:
                    text = ' '.join([page.extract_text() or "" for page in pdf.pages])
                    print(f"Extracted text from {filename}: {text[:100]}...")  # Debug
                    return self.clean_text(text)
            
            # DOCX files
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or file_ext == '.docx':
                try:
                    import docx
                except ImportError:
                    return f"Error extracting text from {filename}: python-docx module not installed"
                doc = docx.Document(file_path)
                text = ' '.join([para.text for para in doc.paragraphs if para.text])
                print(f"Extracted text from {filename}: {text[:100]}...")  # Debug
                return self.clean_text(text)
            
            # PPT/PPTX files
            elif mime_type in ['application/vnd.ms-powerpoint', 
                             'application/vnd.openxmlformats-officedocument.presentationml.presentation'] or file_ext in ['.ppt', '.pptx']:
                try:
                    from pptx import Presentation
                except ImportError:
                    return f"Error extracting text from {filename}: python-pptx module not installed"
                prs = Presentation(file_path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            text_parts.append(shape.text)
                text = ' '.join(text_parts)
                print(f"Extracted text from {filename}: {text[:100]}...")  # Debug
                return self.clean_text(text)
            
            # Excel files (XLSX/XLS)
            elif mime_type in ['application/vnd.ms-excel', 
                             'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'] or file_ext in ['.xlsx', '.xls']:
                try:
                    import pandas as pd
                except ImportError:
                    return f"Error extracting text from {filename}: pandas module not installed"
                try:
                    xls = pd.ExcelFile(file_path)
                    text_parts = []
                    for sheet_name in xls.sheet_names:
                        df = pd.read_excel(file_path, sheet_name=sheet_name)
                        text_parts.append(f"Sheet: {sheet_name}")
                        text_parts.append(str(df))  # Extract all rows
                    text = "\n".join(text_parts)
                    print(f"Extracted text from {filename}: {text[:100]}...")  # Debug
                    return self.clean_text(text)
                except Exception as e:
                    return f"Error extracting text from {filename}: {str(e)}"
            
            # JPEG and PNG images
            elif mime_type in ['image/jpeg', 'image/png'] or file_ext in ['.jpg', '.jpeg', '.png']:
                try:
                    from PIL import Image
                    import pytesseract
                    import cv2
                except ImportError:
                    return f"Error extracting text from {filename}: PIL, pytesseract, or opencv-python module not installed"
                # Preprocess image
                img = cv2.imread(file_path)
                gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
                _, thresh = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY)
                temp_img_path = f"temp_{filename}.png"
                cv2.imwrite(temp_img_path, thresh)
                text = pytesseract.image_to_string(Image.open(temp_img_path))
                os.remove(temp_img_path)
                print(f"Extracted text from {filename}: {text[:100]}...")  # Debug
                return self.clean_text(text)
            
            return f"File type {mime_type or file_ext} not supported for text extraction"
        
        except Exception as e:
            return f"Error extracting text from {filename}: {str(e)}"
    
    def generate_file_embeddings(self):
        """
        Scan directory, generate embeddings for valid files and store in vector database
        """
        if not os.path.exists(self.base_path):
            print(f"Directory {self.base_path} does not exist. No files processed.")
            return
        
        # Reset FAISS index and metadata
        self.index = faiss.IndexFlatL2(self.dimension)
        self.file_metadata = {}
        
        # Collect texts and metadata for batch processing
        file_texts = []
        file_hashes = []
        file_infos = []
        
        for root, _, files in os.walk(self.base_path):
            for filename in files:
                file_path = os.path.join(root, filename)
                
                if self.is_valid_file(file_path):
                    try:
                        # Generate unique hash and metadata
                        file_hash = hashlib.md5(file_path.encode()).hexdigest()
                        mod_time = os.path.getmtime(file_path)
                        
                        # Extract and chunk text
                        file_text = self.extract_file_text(file_path)
                        if file_text.startswith("Error extracting"):
                            print(file_text)
                            continue
                        
                        # Chunk text to avoid truncation
                        text_chunks = self.chunk_text(file_text)
                        if not text_chunks:
                            print(f"Skipping {filename}: No valid text chunks")
                            continue
                        
                        # Prefix chunks with 'passage:' for e5-large-v2
                        prefixed_chunks = [f"passage: {chunk}" for chunk in text_chunks]
                        
                        # Store for batch processing
                        file_texts.extend(prefixed_chunks)
                        file_hashes.extend([file_hash] * len(prefixed_chunks))
                        file_infos.extend([{
                            'path': file_path,
                            'filename': filename,
                            'size': os.path.getsize(file_path),
                            'mod_time': mod_time
                        }] * len(prefixed_chunks))
                        print(f"Queued {filename} with {len(prefixed_chunks)} chunks")
                    
                    except Exception as e:
                        print(f"Error processing {file_path}: {e}")
        
        # Batch encode embeddings
        if file_texts:
            try:
                embeddings = self.model.encode(file_texts, batch_size=16, show_progress_bar=True, normalize_embeddings=True)
                # Normalize embeddings (redundant but ensures consistency)
                embeddings = normalize(embeddings)
                self.index.add(embeddings)
                
                # Store metadata
                for file_hash, file_info in zip(file_hashes, file_infos):
                    if file_hash not in self.file_metadata:
                        self.file_metadata[file_hash] = file_info
                print(f"Processed {len(file_texts)} text chunks from {len(self.file_metadata)} files")
            except Exception as e:
                print(f"Error encoding embeddings: {e}")
    
    def recommend_files(self, query: str, top_k: int = 5) -> List[Dict]:
        """
        Recommend files based on text query
        
        Args:
            query (str): User's text query
            top_k (int): Number of top recommendations to return
        
        Returns:
            List of recommended files with metadata
        """
        if not self.file_metadata:
            print("No files indexed. Run generate_file_embeddings first.")
            return []
        
        # Prefix query with 'query:' for e5-large-v2
        prefixed_query = f"query: {query}"
        
        # Generate and normalize query embedding
        query_embedding = self.model.encode(prefixed_query, normalize_embeddings=True)
        query_embedding = normalize([query_embedding])[0]
        
        # Search in vector database
        distances, indices = self.index.search(np.array([query_embedding]), min(top_k, self.index.ntotal))
        
        # Compute cosine similarity
        recommendations = []
        for dist, idx in zip(distances[0], indices[0]):
            file_hash = list(self.file_metadata.keys())[idx % len(self.file_metadata)]
            file_info = self.file_metadata[file_hash].copy()
            # Convert L2 distance to cosine similarity
            l2_dist = dist
            cosine_sim = 1 - (l2_dist / 2)  # Approximate cosine from L2 for normalized vectors
            file_info['similarity_score'] = max(0, min(1, cosine_sim))  # Clamp to [0, 1]
            recommendations.append(file_info)
        
        # Deduplicate by file path
        seen_paths = set()
        unique_recommendations = []
        for rec in recommendations:
            if rec['path'] not in seen_paths:
                unique_recommendations.append(rec)
                seen_paths.add(rec['path'])
        
        return unique_recommendations[:top_k]

# Example usage
if __name__ == "__main__":
    # Initialize the system with base path
    recommender = FileRecommendationSystem("E:\\Documents")
    
    # Generate embeddings for all files
    recommender.generate_file_embeddings()
    
    # Example recommendation
    results = recommender.recommend_files("maven")
    if results:
        print("\nRecommended files:")
        for result in results:
            print(f"Recommended File: {result['filename']}")
            print(f"Path: {result['path']}")
            print(f"Similarity Score: {result['similarity_score']:.4f}\n")
    else:
        print("No recommendations found.")