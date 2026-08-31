import os
import hashlib
import numpy as np
from sentence_transformers import SentenceTransformer
from typing import List, Dict
import faiss
import mimetypes
import re

class FileRecommendationSystem:
    def __init__(self, base_path: str, embedding_model: str = 'all-MiniLM-L6-v2'):
        """
        Initialize the File Recommendation System
        
        Args:
            base_path (str): Root directory to scan for files
            embedding_model (str): Sentence Transformer model for generating embeddings
        """
        self.base_path = base_path
        self.embedding_model = embedding_model
        self.model = SentenceTransformer(embedding_model)
        
        # System and temporary file extensions to exclude
        self.excluded_extensions = {
            '.tmp', '.log', '.swp', '.lock', 
            '.sys', '.dat', '.ini', '.config',
            '.exe', '.dll', '.bin', 
            '.DS_Store', 'desktop.ini', 'thumbs.db'
        }
        
        # Supported file extensions
        self.supported_extensions = {'.txt', '.pdf', '.docx', '.ppt', '.pptx', '.jpg', '.jpeg', '.xlsx', '.xls', '.png', '.cpp'}
        
        # Vector database using FAISS
        self.dimension = self.model.get_sentence_embedding_dimension()
        self.index = faiss.IndexFlatL2(self.dimension)
        
        # Metadata storage
        self.file_metadata = {}

    def is_valid_file(self, file_path: str) -> bool:
        """
        Check if file should be included in recommendation system
        
        Args:
            file_path (str): Path to the file
        
        Returns:
            bool: Whether file is valid for processing
        """
        filename = os.path.basename(file_path)
        file_ext = os.path.splitext(filename)[1].lower()
        
        conditions = [
            not filename.startswith('.'),  # Exclude hidden files
            file_ext not in self.excluded_extensions,
            file_ext in self.supported_extensions,  # Only allow supported files
            os.path.isfile(file_path),  # Ensure it's a file
            os.path.getsize(file_path) > 0  # Exclude empty files
        ]
        
        return all(conditions)
    
    def clean_text(self, text: str) -> str:
        """
        Clean extracted text by removing excessive whitespace and special characters
        
        Args:
            text (str): Raw extracted text
        
        Returns:
            str: Cleaned text
        """
        text = re.sub(r'\s+', ' ', text.strip())
        return text if text else "No text extracted"
    
    def extract_file_text(self, file_path: str) -> str:
        """
        Extract text content from different file types
        
        Args:
            file_path (str): Path to the file
        
        Returns:
            str: Extracted text content or error message
        """
        try:
            mime_type, _ = mimetypes.guess_type(file_path)
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # Text files
            if mime_type and mime_type.startswith('text/'):
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return self.clean_text(f.read())
            
            # PDF files
            elif mime_type == 'application/pdf' or file_ext == '.pdf':
                try:
                    import PyPDF2
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: PyPDF2 module not installed"
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ' '.join([page.extract_text() or "" for page in reader.pages])
                    return self.clean_text(text)
            
            # DOCX files
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or file_ext == '.docx':
                try:
                    import docx
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: python-docx module not installed"
                doc = docx.Document(file_path)
                text = ' '.join([para.text for para in doc.paragraphs if para.text])
                return self.clean_text(text)
            
            # PPT/PPTX files
            elif mime_type in ['application/vnd.ms-powerpoint', 
                             'application/vnd.openxmlformats-officedocument.presentationml.presentation'] or file_ext in ['.ppt', '.pptx']:
                try:
                    from pptx import Presentation
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: python-pptx module not installed"
                prs = Presentation(file_path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            text_parts.append(shape.text)
                text = ' '.join(text_parts)
                return self.clean_text(text)
            
            # Excel files (XLSX/XLS)
            elif mime_type in ['application/vnd.ms-excel', 
                             'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'] or file_ext in ['.xlsx', '.xls']:
                try:
                    import pandas as pd
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: pandas module not installed"
                try:
                    xls = pd.ExcelFile(file_path)
                    text_parts = []
                    for sheet_name in xls.sheet_names:
                        df = pd.read_excel(file_path, sheet_name=sheet_name)
                        text_parts.append(f"Sheet: {sheet_name}")
                        text_parts.append(str(df.head(100)))  # Limit to first 100 rows
                    text = "\n".join(text_parts)
                    return self.clean_text(text)
                except Exception as e:
                    return f"Error extracting text from {os.path.basename(file_path)}: {str(e)}"
            
            # JPEG and PNG images
            elif mime_type in ['image/jpeg', 'image/png'] or file_ext in ['.jpg', '.jpeg', '.png']:
                try:
                    import cv2
                    import easyocr
                except ImportError:
                    return f"Error extracting text from {os.path.basename(file_path)}: cv2 or easyocr module not installed"
                image = cv2.imread(file_path)
                if image is None:
                    return f"Error reading image {os.path.basename(file_path)}: Unable to load image"
                reader = easyocr.Reader(['en'])  # Initialize EasyOCR for English
                text_parts = reader.readtext(image, detail=0)  # Extract text without bounding box details
                text = ' '.join(text_parts)
                return self.clean_text(text)
            
            return f"File type {mime_type or file_ext} not supported for text extraction"
        
        except Exception as e:
            return f"Error extracting text from {os.path.basename(file_path)}: {str(e)}"
    
    def generate_file_embeddings(self):
        """
        Scan directory, generate embeddings for valid files and store in vector database
        """
        if not os.path.exists(self.base_path):
            print(f"Directory {self.base_path} does not exist. No files processed.")
            return
        
        # Reset FAISS index and metadata to ensure fresh processing
        self.index = faiss.IndexFlatL2(self.dimension)
        self.file_metadata = {}
        
        for root, _, files in os.walk(self.base_path):
            for filename in files:
                file_path = os.path.join(root, filename)
                
                if self.is_valid_file(file_path):
                    try:
                        # Generate a unique hash for the file
                        file_hash = hashlib.md5(file_path.encode()).hexdigest()
                        mod_time = os.path.getmtime(file_path)
                        
                        # Extract text content
                        file_text = self.extract_file_text(file_path)
                        
                        # Skip files with extraction errors
                        if file_text.startswith("Error extracting"):
                            print(file_text)
                            continue
                        
                        # Generate embedding
                        embedding = self.model.encode(file_text)
                        
                        # Add to FAISS index
                        self.index.add(np.array([embedding]))
                        
                        # Store metadata
                        self.file_metadata[file_hash] = {
                            'path': file_path,
                            'filename': filename,
                            'size': os.path.getsize(file_path),
                            'mod_time': mod_time
                        }
                        print(f"Processed {filename}")
                    
                    except Exception as e:
                        print(f"Error processing {file_path}: {e}")
    
    def recommend_files(self, query: str, top_k: int = 5) -> List[Dict]:
        """
        Recommend files based on text query
        
        Args:
            query (str): User's text query
            top_k (int): Number of top recommendations to return
        
        Returns:
            List of recommended files with metadata
        """
        if not self.file_metadata:
            print("No files indexed. Run generate_file_embeddings first.")
            return []
        
        # Generate embedding for query
        query_embedding = self.model.encode(query)
        
        # Search in vector database
        distances, indices = self.index.search(np.array([query_embedding]), min(top_k, len(self.file_metadata)))
        
        # Retrieve and return recommended files
        recommendations = []
        for dist, idx in zip(distances[0], indices[0]):
            file_hash = list(self.file_metadata.keys())[idx]
            file_info = self.file_metadata[file_hash].copy()
            file_info['similarity_score'] = 1 / (1 + dist)  # Convert distance to similarity
            recommendations.append(file_info)
        
        return recommendations

# Example usage
if __name__ == "__main__":
    # Initialize the system with base path
    recommender = FileRecommendationSystem("E:\\Documents")
    
    # Generate embeddings for all files
    recommender.generate_file_embeddings()
    
    # Example recommendation
    results = recommender.recommend_files("cuckoo solution")
    if results:
        print("\nRecommended files:")
        for result in results:
            print(f"Recommended File: {result['filename']}")
            print(f"Path: {result['path']}")
            print(f"Similarity Score: {result['similarity_score']:.4f}\n")
    else:
        print("No recommendations found.")